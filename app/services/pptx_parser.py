"""Parse PPTX and compute lightweight presentation metrics.

Exports parse_pptx_metrics() which returns:
- content_for_llm: slide title/bullets/notes for each slide
- deck_metrics: density, small fonts, contrast, style consistency, VBA flag
"""

from typing import Any, Dict, List, Optional, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _srgb_channel_to_linear(channel: float) -> float:
    if channel <= 0.04045:
        return channel / 12.92
    return ((channel + 0.055) / 1.055) ** 2.4


def _relative_luminance(r: int, g: int, b: int) -> float:
    rl = _srgb_channel_to_linear(r / 255.0)
    gl = _srgb_channel_to_linear(g / 255.0)
    bl = _srgb_channel_to_linear(b / 255.0)
    return 0.2126 * rl + 0.7152 * gl + 0.0722 * bl


def _contrast_ratio(rgb1: Tuple[int, int, int], rgb2: Tuple[int, int, int]) -> float:
    l1 = _relative_luminance(*rgb1)
    l2 = _relative_luminance(*rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _pptx_color_to_rgb(color) -> Optional[Tuple[int, int, int]]:
    try:
        if not color:
            return None
        if color.rgb is not None:
            rgb = color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None
    return None


def _get_slide_background_rgb(slide) -> Optional[Tuple[int, int, int]]:
    try:
        fill = slide.background.fill
        if fill and fill.fore_color:
            return _pptx_color_to_rgb(fill.fore_color)
    except Exception:
        return None
    return None


def parse_pptx_metrics(ppt_path: Path) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
    """Extract slide texts and compute deck metrics from a PPTX/PPTM file.

    Returns (content_for_llm, deck_metrics).
    content_for_llm: list of {index, title, bullets[], notes}
    deck_metrics: density, small_fonts, contrast_issues, style_inconsistency, vba_summary
    """
    prs = Presentation(str(ppt_path))
    slide_size = (prs.slide_width, prs.slide_height)

    slides_summary: List[Dict[str, Any]] = []
    deck_fonts: List[str] = []
    deck_font_sizes: List[float] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_chars = 0
        text_shapes_area = 0
        min_font_pt_on_slide: Optional[float] = None
        font_families_on_slide: List[str] = []
        font_sizes_on_slide: List[float] = []
        contrast_issue = False

        bg_rgb = _get_slide_background_rgb(slide) or (255, 255, 255)

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    if hasattr(subshape, "text_frame") and subshape.text_frame:
                        t = subshape.text_frame.text or ""
                        if not t.strip():
                            continue
                        slide_chars += len(t)
                        try:
                            text_shapes_area += int(subshape.width) * int(subshape.height)
                        except Exception:
                            pass
                        for paragraph in subshape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                size = run.font.size.pt if run.font.size is not None else None
                                if size is not None:
                                    font_sizes_on_slide.append(size)
                                    deck_font_sizes.append(size)
                                    if min_font_pt_on_slide is None or size < min_font_pt_on_slide:
                                        min_font_pt_on_slide = size
                                if run.font.name:
                                    font_families_on_slide.append(run.font.name)
                                    deck_fonts.append(run.font.name)
                                try:
                                    frgb = _pptx_color_to_rgb(run.font.color)
                                    if frgb is not None:
                                        ratio = _contrast_ratio(frgb, bg_rgb)
                                        if ratio < 4.5:
                                            contrast_issue = True
                                except Exception:
                                    pass
            elif hasattr(shape, "text_frame") and shape.text_frame:
                t = shape.text_frame.text or ""
                if not t.strip():
                    continue
                slide_chars += len(t)
                try:
                    text_shapes_area += int(shape.width) * int(shape.height)
                except Exception:
                    pass
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        size = run.font.size.pt if run.font.size is not None else None
                        if size is not None:
                            font_sizes_on_slide.append(size)
                            deck_font_sizes.append(size)
                            if min_font_pt_on_slide is None or size < min_font_pt_on_slide:
                                min_font_pt_on_slide = size
                        if run.font.name:
                            font_families_on_slide.append(run.font.name)
                            deck_fonts.append(run.font.name)
                        try:
                            frgb = _pptx_color_to_rgb(run.font.color)
                            if frgb is not None:
                                ratio = _contrast_ratio(frgb, bg_rgb)
                                if ratio < 4.5:
                                    contrast_issue = True
                        except Exception:
                            pass

        slide_area = int(slide_size[0]) * int(slide_size[1]) or 1
        text_area_ratio = (text_shapes_area / slide_area) if slide_area else 0.0

        char_density = min(1.0, slide_chars / 900.0)
        density_score = round(0.7 * char_density + 0.3 * min(1.0, text_area_ratio * 3.0), 3)

        slides_summary.append({
            "index": slide_index,
            "chars": slide_chars,
            "text_area_ratio": round(text_area_ratio, 3),
            "density_score": density_score,
            "min_font_pt": min_font_pt_on_slide if min_font_pt_on_slide is not None else None,
            "font_families": list(sorted(set(font_families_on_slide)))[:6],
            "contrast_issue": contrast_issue,
        })

    small_fonts_slides = [s["index"] for s in slides_summary if (s["min_font_pt"] or 1000) < 18.0]
    contrast_issues_slides = [s["index"] for s in slides_summary if s["contrast_issue"]]
    text_density_avg = round(sum(s["density_score"] for s in slides_summary) / max(1, len(slides_summary)), 3)
    density_bad_on = [s["index"] for s in slides_summary if s["density_score"] > 0.7]

    majority_font = None
    if deck_fonts:
        from collections import Counter
        majority_font = Counter(deck_fonts).most_common(1)[0][0]
    avg_size = None
    if deck_font_sizes:
        avg_size = round(sum(deck_font_sizes) / len(deck_font_sizes), 1)

    style_inconsistency_slides: List[int] = []
    if majority_font is not None and avg_size is not None:
        for s in slides_summary:
            if (majority_font not in s["font_families"]) or (s["min_font_pt"] is not None and abs(s["min_font_pt"] - avg_size) >= 6):
                style_inconsistency_slides.append(s["index"])

    vba_summary = {"has_vba": ppt_path.suffix.lower().endswith("pptm"), "note": "macros not parsed in MVP"}

    deck_metrics = {
        "text_density": {"avg": text_density_avg, "bad_on": density_bad_on},
        "small_fonts": small_fonts_slides,
        "contrast_issues": contrast_issues_slides,
        "style_inconsistency": style_inconsistency_slides,
        "vba_summary": vba_summary,
    }

    content_for_llm: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i}"
        bullets: List[str] = []
        notes = ""
        for shp in slide.shapes:
            if hasattr(shp, "text_frame") and shp.text_frame:
                tt = (shp.text_frame.text or "").strip()
                if tt:
                    bullets.append(tt)
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        content_for_llm.append({
            "index": i,
            "title": title,
            "bullets": [b for b in bullets if b],
            "notes": notes,
        })

    return content_for_llm, deck_metrics


